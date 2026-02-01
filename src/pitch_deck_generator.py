"""
Executive pitch deck generator with Apple-inspired aesthetics.

Creates a beautiful, minimal, high-impact presentation for C-suite executives.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import json


class AppleStylePitchDeck:
    """
    Creates an Apple-inspired pitch deck with clean aesthetics.

    Design principles:
    - Abundant white space
    - Minimal text, maximum impact
    - Clean sans-serif typography
    - Simple, elegant color palette
    - High-quality visuals
    """

    # Apple-inspired color palette
    COLORS = {
        'black': RGBColor(0, 0, 0),
        'dark_gray': RGBColor(51, 51, 51),
        'gray': RGBColor(142, 142, 147),
        'light_gray': RGBColor(229, 229, 234),
        'white': RGBColor(255, 255, 255),
        'blue': RGBColor(0, 122, 255),
        'green': RGBColor(52, 199, 89),
        'orange': RGBColor(255, 149, 0),
        'red': RGBColor(255, 59, 48),
        'subtle_gray': RGBColor(242, 242, 247)
    }

    def __init__(self, results_path: str = "outputs/evaluation_results.json"):
        """Initialize with evaluation results."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        with open(results_path, 'r') as f:
            self.results = json.load(f)

        self.figures_dir = Path("outputs/figures")

    def _add_title_slide(self):
        """Create a minimal, impactful title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Main title - centered, large
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5), Inches(12), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.text = "Evaluating Human Simulations"
        p.font.size = Pt(66)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['black']
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(5.2), Inches(12), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame

        p = subtitle_frame.paragraphs[0]
        p.text = "A Multi-Dimensional Assessment of AI-Generated Consumer Personas"
        p.font.size = Pt(28)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # Thin line accent
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(6), Inches(6.5), Inches(4), Inches(0)
        )
        line.line.color.rgb = self.COLORS['blue']
        line.line.width = Pt(3)

    def _add_research_slide(self):
        """Slide 1: Research Approach & Methodology."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8), Inches(14), Inches(1)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Our Approach"
        p.font.size = Pt(54)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['black']
        p.font.bold = True

        # Subtitle - The Challenge
        challenge_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8), Inches(7), Inches(0.6)
        )
        p = challenge_box.text_frame.paragraphs[0]
        p.text = "The Challenge"
        p.font.size = Pt(24)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['blue']
        p.font.bold = True

        # Challenge description
        desc_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(7), Inches(1.5)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = "How accurately do our AI-generated human simulations capture the nuanced behaviors, language patterns, and personality traits of real consumers?"
        p.font.size = Pt(20)
        p.font.name = 'SF Pro Text'
        p.font.color.rgb = self.COLORS['dark_gray']
        p.line_spacing = 1.4

        # Our Solution title
        solution_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.3), Inches(7), Inches(0.6)
        )
        p = solution_box.text_frame.paragraphs[0]
        p.text = "Our Solution"
        p.font.size = Pt(24)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['green']
        p.font.bold = True

        # Three-pillar approach
        pillars_y = 5.0
        pillar_height = 1.2

        pillars = [
            ("Semantic Similarity", "What they say", "Measures content alignment using state-of-the-art embeddings"),
            ("Stylistic Analysis", "How they say it", "Captures writing style, tone, and linguistic patterns"),
            ("Human Benchmarking", "Gold standard", "Compares against 85% human test-retest reliability")
        ]

        for i, (title, subtitle, desc) in enumerate(pillars):
            y_pos = pillars_y + (i * 1.15)

            # Number circle
            circle = slide.shapes.add_shape(
                5,  # Oval
                Inches(1), Inches(y_pos + 0.05), Inches(0.35), Inches(0.35)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.COLORS['blue']
            circle.line.fill.background()

            # Number text
            num_frame = circle.text_frame
            num_p = num_frame.paragraphs[0]
            num_p.text = str(i + 1)
            num_p.font.size = Pt(16)
            num_p.font.color.rgb = self.COLORS['white']
            num_p.font.bold = True
            num_p.alignment = PP_ALIGN.CENTER
            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Title
            t_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_pos), Inches(6), Inches(0.4)
            )
            t_p = t_box.text_frame.paragraphs[0]
            t_p.text = f"{title} • {subtitle}"
            t_p.font.size = Pt(18)
            t_p.font.name = 'SF Pro Display'
            t_p.font.color.rgb = self.COLORS['black']
            t_p.font.bold = True

            # Description
            d_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_pos + 0.35), Inches(6), Inches(0.5)
            )
            d_frame = d_box.text_frame
            d_frame.word_wrap = True
            d_p = d_frame.paragraphs[0]
            d_p.text = desc
            d_p.font.size = Pt(14)
            d_p.font.name = 'SF Pro Text'
            d_p.font.color.rgb = self.COLORS['gray']

        # Add visualization on right
        if (self.figures_dir / "overview_scores.png").exists():
            slide.shapes.add_picture(
                str(self.figures_dir / "overview_scores.png"),
                Inches(9), Inches(2.3), height=Inches(5)
            )

    def _add_results_slide(self):
        """Slide 2: Key Results."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8), Inches(14), Inches(1)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = "What We Found"
        p.font.size = Pt(54)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['black']
        p.font.bold = True

        # Get results
        agg = self.results['aggregate_scores']

        # Big number callout - Performance Gap
        gap_percentage = ((0.85 - agg['semantic']['mean_cosine_similarity']) / 0.85) * 100

        big_num_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(6), Inches(2)
        )
        big_frame = big_num_box.text_frame

        # Number
        p = big_frame.paragraphs[0]
        p.text = f"{gap_percentage:.0f}%"
        p.font.size = Pt(120)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['red']
        p.font.bold = True

        # Label below number
        p = big_frame.add_paragraph()
        p.text = "Below Human Benchmark"
        p.font.size = Pt(28)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['dark_gray']
        p.space_before = Pt(0)

        p = big_frame.add_paragraph()
        p.text = "Our simulations achieved 52% accuracy vs. 85% human test-retest reliability"
        p.font.size = Pt(16)
        p.font.name = 'SF Pro Text'
        p.font.color.rgb = self.COLORS['gray']
        p.space_before = Pt(10)

        # Key metrics boxes
        metrics_y = 4.5
        metric_data = [
            ("2.82×", "Response Length", "AI responses are nearly 3× longer than human responses", self.COLORS['orange']),
            ("0%", "Natural Imperfections", "AI responses lack typos, casual language, and human quirks", self.COLORS['red']),
            ("2.66/10", "Worst Category", "Abstract questions (influences) show critical weakness", self.COLORS['orange'])
        ]

        x_positions = [1, 5.5, 10]
        for (num, label, desc, color), x_pos in zip(metric_data, x_positions):
            # Background box
            box = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(x_pos), Inches(metrics_y), Inches(4), Inches(2.2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS['subtle_gray']
            box.line.fill.background()

            # Number
            num_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.3), Inches(metrics_y + 0.3), Inches(3.4), Inches(0.8)
            )
            p = num_box.text_frame.paragraphs[0]
            p.text = num
            p.font.size = Pt(48)
            p.font.name = 'SF Pro Display'
            p.font.color.rgb = color
            p.font.bold = True

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.3), Inches(metrics_y + 1.1), Inches(3.4), Inches(0.35)
            )
            p = label_box.text_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(20)
            p.font.name = 'SF Pro Display'
            p.font.color.rgb = self.COLORS['black']
            p.font.bold = True

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.3), Inches(metrics_y + 1.5), Inches(3.4), Inches(0.6)
            )
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            p = desc_frame.paragraphs[0]
            p.text = desc
            p.font.size = Pt(13)
            p.font.name = 'SF Pro Text'
            p.font.color.rgb = self.COLORS['gray']
            p.line_spacing = 1.3

        # Add chart
        if (self.figures_dir / "dimension_comparison.png").exists():
            slide.shapes.add_picture(
                str(self.figures_dir / "dimension_comparison.png"),
                Inches(8), Inches(1.8), width=Inches(7)
            )

    def _add_insights_slide(self):
        """Slide 3: Insights & Recommendations."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8), Inches(14), Inches(1)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = "Path Forward"
        p.font.size = Pt(54)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['black']
        p.font.bold = True

        # Left column - Insights
        insight_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(7), Inches(0.5)
        )
        p = insight_box.text_frame.paragraphs[0]
        p.text = "Key Insights"
        p.font.size = Pt(28)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['blue']
        p.font.bold = True

        insights_text = slide.shapes.add_textbox(
            Inches(1), Inches(2.7), Inches(7), Inches(4)
        )
        frame = insights_text.text_frame
        frame.word_wrap = True

        insights = [
            "The simulation captures semantic content moderately well (52%) but fails to replicate human behavioral authenticity",
            "Length mismatch is the most obvious tell — AI responses are nearly 3× more verbose than human responses",
            "Complete absence of natural imperfections makes responses detectably artificial",
            "Performance varies significantly by individual (47-55%) and question type (27-59%)",
            "Abstract questions reveal fundamental limitations in personality simulation"
        ]

        for i, insight in enumerate(insights):
            if i > 0:
                p = frame.add_paragraph()
            else:
                p = frame.paragraphs[0]

            p.text = f"•  {insight}"
            p.font.size = Pt(16)
            p.font.name = 'SF Pro Text'
            p.font.color.rgb = self.COLORS['dark_gray']
            p.line_spacing = 1.5
            p.space_before = Pt(12) if i > 0 else Pt(0)

        # Right column - Recommendations
        rec_box = slide.shapes.add_textbox(
            Inches(9), Inches(2), Inches(6), Inches(0.5)
        )
        p = rec_box.text_frame.paragraphs[0]
        p.text = "Immediate Actions"
        p.font.size = Pt(28)
        p.font.name = 'SF Pro Display'
        p.font.color.rgb = self.COLORS['green']
        p.font.bold = True

        # Recommendation cards
        recs = [
            ("1", "Calibrate Response Length", "Constrain AI verbosity to match human patterns (target: 1.0-1.2× vs current 2.82×)"),
            ("2", "Inject Naturalness", "Add controlled imperfections: casual language, contractions, occasional typos"),
            ("3", "Refine Personality Capture", "Individual-specific tuning, especially for underperforming personas")
        ]

        rec_y = 2.8
        for num, title, desc in recs:
            # Card background
            card = slide.shapes.add_shape(
                1,
                Inches(9), Inches(rec_y), Inches(6), Inches(1.3)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = self.COLORS['light_gray']
            card.line.width = Pt(1)

            # Number badge
            badge = slide.shapes.add_shape(
                5,  # Oval
                Inches(9.3), Inches(rec_y + 0.15), Inches(0.4), Inches(0.4)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.COLORS['green']
            badge.line.fill.background()

            badge_p = badge.text_frame.paragraphs[0]
            badge_p.text = num
            badge_p.font.size = Pt(18)
            badge_p.font.color.rgb = self.COLORS['white']
            badge_p.font.bold = True
            badge_p.alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Title
            t_box = slide.shapes.add_textbox(
                Inches(9.85), Inches(rec_y + 0.2), Inches(4.8), Inches(0.35)
            )
            t_p = t_box.text_frame.paragraphs[0]
            t_p.text = title
            t_p.font.size = Pt(18)
            t_p.font.name = 'SF Pro Display'
            t_p.font.color.rgb = self.COLORS['black']
            t_p.font.bold = True

            # Description
            d_box = slide.shapes.add_textbox(
                Inches(9.85), Inches(rec_y + 0.55), Inches(4.8), Inches(0.6)
            )
            d_frame = d_box.text_frame
            d_frame.word_wrap = True
            d_p = d_frame.paragraphs[0]
            d_p.text = desc
            d_p.font.size = Pt(13)
            d_p.font.name = 'SF Pro Text'
            d_p.font.color.rgb = self.COLORS['gray']
            d_p.line_spacing = 1.3

            rec_y += 1.5

        # Bottom insight box
        bottom_box = slide.shapes.add_shape(
            1,
            Inches(1), Inches(7.2), Inches(14), Inches(1.2)
        )
        bottom_box.fill.solid()
        bottom_box.fill.fore_color.rgb = self.COLORS['blue']
        bottom_box.line.fill.background()

        conclusion_text = slide.shapes.add_textbox(
            Inches(1.5), Inches(7.4), Inches(13), Inches(0.8)
        )
        c_frame = conclusion_text.text_frame
        c_frame.word_wrap = True

        c_p = c_frame.paragraphs[0]
        c_p.text = "Bottom Line: "
        run = c_p.runs[0]
        run.font.size = Pt(20)
        run.font.name = 'SF Pro Display'
        run.font.color.rgb = self.COLORS['white']
        run.font.bold = True

        run = c_p.add_run()
        run.text = "The simulations require significant refinement before production deployment. With targeted improvements to length calibration and naturalness, we can close the 38% performance gap."
        run.font.size = Pt(20)
        run.font.name = 'SF Pro Text'
        run.font.color.rgb = self.COLORS['white']

        c_p.line_spacing = 1.4

    def generate(self, output_path: str = "outputs/executive_pitch_deck.pptx"):
        """Generate the complete pitch deck."""
        print("Creating Apple-style executive pitch deck...")

        self._add_title_slide()
        self._add_research_slide()
        self._add_results_slide()
        self._add_insights_slide()

        # Save
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_file))

        print(f"✅ Pitch deck created: {output_path}")
        return output_path


if __name__ == "__main__":
    deck = AppleStylePitchDeck()
    deck.generate()
    print("\n🎯 Executive pitch deck ready for CEO presentation!")
